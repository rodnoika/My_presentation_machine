import openai
from pptx import Presentation
from pptx.util import Inches

openai.api_key = 'API'

def generate_slide_content(prompt):
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo-1106", 
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ]
    )
    return response['choices'][0]['message']['content'].strip()


def create_presentation(prompt, num_slides=10):
    prs = Presentation()

    # Генерация содержания слайдов
    for i in range(num_slides):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        slide_prompt = f"{prompt} - Слайд {i+1}:"
        slide_content = generate_slide_content(slide_prompt)
        
        title.text = f"Слайд {i+1}"
        content.text = slide_content

    prs.save('presentation.pptx')
    print("Презентация создана и сохранена")

prompt = input()
create_presentation(prompt)
